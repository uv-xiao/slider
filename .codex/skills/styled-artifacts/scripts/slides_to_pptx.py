#!/usr/bin/env python3
"""
Combine slide images into a single PPTX (each slide is the image).

This is a pragmatic bridge: generate high-quality slide images (e.g., via Nano Banana Pro),
then package them into a PowerPoint deck for distribution/edits.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import List

try:
    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.util import Inches  # type: ignore[import-not-found]
except ImportError:
    print("Error: python-pptx not found. Install with: pip install python-pptx")
    raise SystemExit(1)


def get_image_files(paths: List[str]) -> List[Path]:
    image_extensions = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
    image_files: list[Path] = []

    for path_str in paths:
        path = Path(path_str)
        if path.is_file():
            if path.suffix.lower() in image_extensions:
                image_files.append(path)
            else:
                print(f"Warning: Skipping non-image file: {path}")
        elif path.is_dir():
            for ext in image_extensions:
                image_files.extend(path.glob(f"*{ext}"))
                image_files.extend(path.glob(f"*{ext.upper()}"))
        else:
            parent = path.parent
            pattern = path.name
            if parent.exists():
                for match in parent.glob(pattern):
                    if match.suffix.lower() in image_extensions:
                        image_files.append(match)

    # De-dup and order by filename.
    image_files = sorted(set(image_files), key=lambda p: p.name)
    return image_files


def _fit_contain(img_w: int, img_h: int, box_w: int, box_h: int) -> tuple[int, int, int, int]:
    """
    Compute a contain-fit rectangle for an image inside a box.
    Returns (left, top, width, height) in the same units.
    """
    if img_w <= 0 or img_h <= 0:
        return (0, 0, box_w, box_h)
    scale = min(box_w / img_w, box_h / img_h)
    w = int(img_w * scale)
    h = int(img_h * scale)
    left = (box_w - w) // 2
    top = (box_h - h) // 2
    return left, top, w, h


def combine_images_to_pptx(image_paths: List[Path], output_path: Path, verbose: bool = False) -> bool:
    if not image_paths:
        print("Error: No image files found")
        return False

    try:
        from PIL import Image  # type: ignore[import-not-found]
    except ImportError:
        print("Error: Pillow not found (needed to size images). Install with: pip install Pillow")
        return False

    prs = Presentation()
    # Default to 16:9 widescreen.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    blank = prs.slide_layouts[6]

    for i, img_path in enumerate(image_paths, start=1):
        slide = prs.slides.add_slide(blank)
        try:
            with Image.open(img_path) as img:
                img_w, img_h = img.size
        except Exception:
            img_w, img_h = (slide_w, slide_h)

        left, top, width, height = _fit_contain(img_w, img_h, slide_w, slide_h)
        slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)

        if verbose:
            print(f"  [{i}/{len(image_paths)}] Added: {img_path.name}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return True


def main() -> int:
    parser = argparse.ArgumentParser(description="Combine slide images into a PPTX presentation")
    parser.add_argument("images", nargs="+", help="Image files, directories, or glob patterns")
    parser.add_argument("-o", "--output", required=True, help="Output PPTX file path")
    parser.add_argument("-v", "--verbose", action="store_true", help="Verbose output")
    args = parser.parse_args()

    image_files = get_image_files(args.images)
    if not image_files:
        print("Error: No image files found matching the specified paths")
        return 1

    if args.verbose:
        print(f"Found {len(image_files)} image(s)")

    ok = combine_images_to_pptx(image_files, Path(args.output), verbose=args.verbose)
    if not ok:
        return 1
    print(str(Path(args.output).resolve()))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

